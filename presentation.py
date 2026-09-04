import os
from collections import defaultdict
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from bond import Bond

# --- MAPPINGS ---
BOND_TYPE_MAP = {
    "GO": "General Obligation",
    "REV": "Revenue Bonds",
    "COP": "Certificate of Participation",
}

# Explicitly map known Bloomberg purposes to the specific slices you want
INCOME_SOURCE_MAP = {
    "SCHOOL IMPS.": "Ad Valorem Property Tax",
    "UNIV. & COLLEGE IMPS.": "College & Univ. Rev.",
    "WATER UTILITY IMPS.": "Water Revenue",
    "PRT, AIRPRT & MARINA IMPS": "Prt, Airprt & Marina Rev.",
    "Local GO": "Ad Valorem Property Tax",
    "State GO": "State General Fund",
    "Higher Education": "College & Univ. Rev.",
    "Airport": "Prt, Airprt & Marina Rev.",
    "Port": "Prt, Airprt & Marina Rev.",
    "Marina": "Prt, Airprt & Marina Rev."
}

CORE_INCOME_BUCKETS = {
    "Ad Valorem Property Tax",
    "State General Fund",
    "Water Revenue",
    "Sewer Revenue",
    "Water & Sewer Revenue",
    "Prt, Airprt & Marina Rev.",
    "College & Univ. Rev."
}

class PresentationGenerator:
    def __init__(self, template_path, output_path, bonds):
        self.template_path = template_path
        self.output_path = output_path
        self.bonds = bonds
        
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template not found: {self.template_path}")
            
        self.prs = Presentation(self.template_path)
    
    def generate(self):
        for i, slide in enumerate(self.prs.slides):
            slide_number = i + 1
            if slide_number == 3:
                self.populate_slide_3(slide)
            else:
                pass
                
        self.prs.save(self.output_path)
        print(f"\nSuccessfully generated presentation: {self.output_path}")

    def populate_slide_3(self, slide):
        issue_types = defaultdict(float)
        income_sources = defaultdict(float)

        for bond in self.bonds:
            # 1. Map Bond Type
            raw_type = bond.muni_issue_type if bond.muni_issue_type != "Unknown" else "Other"
            
            clean_type = BOND_TYPE_MAP.get(raw_type, raw_type.title())
            if "General Obligation" in clean_type:
                clean_type = "General Obligation"
                
            issue_types[clean_type] += bond.market_value

            # 2. Map Income Source (You can change muni_purpose to industry_subgroup if preferred)
            raw_purpose = bond.muni_purpose if bond.muni_purpose != "Unknown" else "Other"
            
            if raw_purpose == "Water & Sewer":
                name = bond.security_name.upper()
                if "WATER" in name and "SEWER" not in name:
                    clean_purpose = "Water Revenue"
                elif "SEWER" in name and "WATER" not in name:
                    clean_purpose = "Sewer Revenue"
                else:
                    clean_purpose = "Water & Sewer Revenue"
            else:
                clean_purpose = INCOME_SOURCE_MAP.get(raw_purpose, raw_purpose.title())
            
            if clean_purpose not in CORE_INCOME_BUCKETS:
                clean_purpose = "Miscellaneous Revenue"
                
            income_sources[clean_purpose] += bond.market_value

        # --- Generate Chart Data ---
        chart_data_type = CategoryChartData()
        chart_data_type.categories = list(issue_types.keys())
        chart_data_type.add_series('Bond Types', list(issue_types.values()))

        chart_data_source = CategoryChartData()
        chart_data_source.categories = list(income_sources.keys())
        chart_data_source.add_series('Income Sources', list(income_sources.values()))

        # --- INJECT DATA INTO EXISTING TEMPLATE CHARTS ---
        charts = [shape.chart for shape in slide.shapes if shape.has_chart]

        if len(charts) >= 2:
            charts[0].replace_data(chart_data_type)
            charts[1].replace_data(chart_data_source)
            print("Successfully updated the beautifully styled charts on Slide 3!")
        else:
            print(f"Warning: Found {len(charts)} charts on Slide 3. Expected at least 2.")


def load_bonds_from_excel(filepath):
    """
    Reads the output Excel file into a list of Bond objects.
    """
    print(f"Loading data from {filepath}...")
    df = pd.read_excel(filepath)
    bonds = []
    
    for index, row in df.iterrows():
        # Convert the row to a dictionary so the Bond class can parse it like Bloomberg data
        row_dict = row.to_dict()
        
        # Try to find standard columns for CUSIP and Face Value
        # Adjust these string keys if your Excel columns are named differently!
        cusip = row_dict.get("CUSIP", row_dict.get("ID_CUSIP", f"Unknown_{index}"))
        face_value = row_dict.get("Face Value", row_dict.get("Par Amount", 100000))
        
        # Create the Bond object. It will automatically pull MUNI_ISSUE_TYP, MUNI_PURPOSE, etc. 
        # if those exact column headers exist in your Excel sheet.
        bond = Bond(cusip=str(cusip), face_value=face_value, bbg_data=row_dict)
        
        # If your Excel sheet already has a calculated Market Value column, use it!
        if "Market Value" in row_dict:
            bond.market_value = row_dict["Market Value"]
            
        bonds.append(bond)
        
    print(f"Successfully loaded {len(bonds)} bonds from the spreadsheet.")
    return bonds


if __name__ == "__main__":
    
    excel_file = "Par_Portfolio_Output.xlsx"
    
    if os.path.exists(excel_file):
        # Load live data from the Excel sheet!
        live_bonds = load_bonds_from_excel(excel_file)
        
        generator = PresentationGenerator(
            template_path="WorkingTemplate.pptx", 
            output_path="PAR_Report_Output.pptx", 
            bonds=live_bonds
        )
        generator.generate()
    else:
        print(f"Could not find {excel_file}. Please ensure the file is in the same folder.")
